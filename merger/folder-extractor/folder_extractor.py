#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
folder-extractor – Universal Folder to Text Converter

Extracts text from any folder containing mixed file types:
- Text files (direct read)
- PDFs (text extraction)
- Images (OCR)
- Office documents (.docx, .pptx)
- And more...

Best-effort approach to create AI-ready context from arbitrary directories.

Features:
- Multi-format support
- Optional OCR integration
- Configurable via TOML
- Auto-splitting for large outputs
"""

from __future__ import annotations

import sys
import os
import io
import json
import hashlib
from pathlib import Path
from datetime import datetime
from typing import Any, Dict, List, Optional, Tuple

# Try importing tomllib for config
try:
    import tomllib  # Python 3.11+
except Exception:
    tomllib = None  # type: ignore


# ========= Configuration =========

ENCODING = "utf-8"
DEFAULT_MAX_FILE_BYTES = 10 * 1024 * 1024  # 10 MB per file
DEFAULT_MAX_OUTPUT_BYTES = 5 * 1024 * 1024  # 5 MB per output file

# Directories to ignore
IGNORE_DIR_NAMES = {
    ".git", ".hg", ".svn",
    "__pycache__", ".mypy_cache", ".pytest_cache",
    "node_modules", "dist", "build", ".next",
    ".venv", "venv",
    ".idea", ".vscode", ".DS_Store",
    ".cargo", ".gradle", ".ruff_cache", ".cache"
}

# Files to ignore
IGNORE_FILE_SUFFIXES = {
    ".lock", ".log"
}

# Binary extensions (no text extraction attempt)
BINARY_EXTS = {
    ".mp3", ".wav", ".flac", ".ogg", ".m4a", ".aac",
    ".mp4", ".mkv", ".mov", ".avi", ".webm",
    ".zip", ".gz", ".bz2", ".xz", ".7z", ".rar", ".zst",
    ".ttf", ".otf", ".woff", ".woff2",
    ".so", ".dylib", ".dll", ".exe",
    ".db", ".sqlite", ".sqlite3", ".realm", ".mdb", ".pack", ".idx",
    ".psd", ".ai", ".sketch", ".fig",
}

# Image extensions (OCR candidates)
IMAGE_EXTS = {
    ".png", ".jpg", ".jpeg", ".gif", ".webp", ".avif", ".bmp", ".ico",
    ".tif", ".tiff"
}

# Document extensions
DOCUMENT_EXTS = {
    ".pdf", ".docx", ".doc", ".pptx", ".ppt", ".xlsx", ".xls"
}

# Text extensions
TEXT_EXTS = {
    ".txt", ".md", ".rst", ".py", ".js", ".ts", ".jsx", ".tsx",
    ".json", ".yaml", ".yml", ".toml", ".ini", ".cfg", ".conf",
    ".sh", ".bash", ".html", ".htm", ".css", ".scss", ".xml",
    ".c", ".h", ".cpp", ".hpp", ".java", ".go", ".rs", ".rb",
    ".php", ".pl", ".lua", ".sql", ".cs", ".swift", ".kt",
}

# Language mapping for syntax highlighting
LANG_MAP = {
    "py": "python", "js": "javascript", "ts": "typescript",
    "json": "json", "yaml": "yaml", "yml": "yaml", "md": "markdown",
    "html": "html", "css": "css", "sh": "bash", "sql": "sql",
    "c": "c", "cpp": "cpp", "java": "java", "go": "go", "rs": "rust",
    "rb": "ruby", "php": "php", "txt": "text",
}


# ========= Utils =========

def human_size(n: int) -> str:
    """Convert bytes to human-readable format."""
    units = ["B", "KB", "MB", "GB", "TB"]
    f = float(n)
    i = 0
    while f >= 1024 and i < len(units) - 1:
        f /= 1024
        i += 1
    return f"{f:.1f} {units[i]}"


def file_md5(path: Path) -> str:
    """Calculate MD5 hash of file."""
    h = hashlib.md5()
    try:
        with path.open("rb") as f:
            for chunk in iter(lambda: f.read(1 << 16), b""):
                h.update(chunk)
        return h.hexdigest()
    except Exception:
        return ""


def detect_file_type(path: Path) -> str:
    """
    Detect file type: text, pdf, image, office, binary, unknown
    """
    ext = path.suffix.lower()
    
    if ext in TEXT_EXTS:
        return "text"
    elif ext == ".pdf":
        return "pdf"
    elif ext in IMAGE_EXTS:
        return "image"
    elif ext in {".docx", ".pptx", ".xlsx"}:
        return "office"
    elif ext in BINARY_EXTS:
        return "binary"
    else:
        # Try to sniff
        try:
            with path.open("rb") as f:
                chunk = f.read(8192)
            if not chunk:
                return "text"
            if b"\x00" in chunk:
                return "binary"
            # Try decode
            try:
                chunk.decode(ENCODING)
                return "text"
            except UnicodeDecodeError:
                return "binary"
        except Exception:
            return "unknown"


def language_for(path: Path) -> str:
    """Get language identifier for syntax highlighting."""
    ext = path.suffix.lower().lstrip(".")
    return LANG_MAP.get(ext, "")


# ========= Config Loading =========

def load_config() -> Dict[str, Any]:
    """Load configuration from ~/.config/folder-extractor/config.toml"""
    cfg: Dict[str, Any] = {
        "general": {},
        "ocr": {},
        "extraction": {}
    }
    
    cfg_path = Path.home() / ".config" / "folder-extractor" / "config.toml"
    if not cfg_path.exists():
        return cfg
    
    if not tomllib:
        return cfg
    
    try:
        text = cfg_path.read_text(encoding=ENCODING)
        data = tomllib.loads(text)
        if isinstance(data, dict):
            cfg.update(data)
    except Exception:
        pass
    
    return cfg


# ========= Extractors =========

def extract_text(path: Path, max_bytes: int) -> Tuple[Optional[str], str]:
    """
    Extract text from text file with chunking for large files.
    Returns (text, method)
    """
    try:
        size = path.stat().st_size
        if max_bytes > 0 and size > max_bytes:
            # Chunk large files instead of truncating
            chunks = []
            chunk_num = 0
            with path.open("rb") as f:
                while True:
                    content = f.read(max_bytes)
                    if not content:
                        break
                    chunk_num += 1
                    text = content.decode(ENCODING, errors="replace")
                    chunks.append(f"--- Chunk {chunk_num} ---\n{text}")
            
            full_text = "\n\n".join(chunks)
            full_text += f"\n\n(File chunked into {chunk_num} parts, total size: {human_size(size)})"
            return full_text, "text_chunked"
        else:
            text = path.read_text(encoding=ENCODING, errors="replace")
            return text, "text"
    except Exception as e:
        return None, f"error: {e}"


def extract_pdf(path: Path) -> Tuple[Optional[str], str]:
    """
    Extract text from PDF.
    Returns (text, method)
    """
    # Try PyPDF2
    try:
        import PyPDF2
        text_parts = []
        with path.open("rb") as f:
            reader = PyPDF2.PdfReader(f)
            for page in reader.pages:
                try:
                    text = page.extract_text()
                    if text:  # Filter out None values
                        text_parts.append(text)
                except Exception:
                    continue
        if text_parts:
            return "\n\n".join(text_parts), "pdf_pypdf2"
    except ImportError:
        pass
    except Exception as e:
        return None, f"pdf_error: {e}"
    
    # Try pdfplumber
    try:
        import pdfplumber
        text_parts = []
        with pdfplumber.open(path) as pdf:
            for page in pdf.pages:
                try:
                    text = page.extract_text()
                    if text:
                        text_parts.append(text)
                except Exception:
                    continue
        if text_parts:
            return "\n\n".join(text_parts), "pdf_pdfplumber"
    except ImportError:
        pass
    except Exception as e:
        return None, f"pdf_error: {e}"
    
    return None, "pdf_no_library"


def extract_image_ocr(path: Path, ocr_config: Dict[str, Any]) -> Tuple[Optional[str], str]:
    """
    Extract text from image via OCR.
    Returns (text, method)
    """
    backend = ocr_config.get("backend", "none")
    
    if backend == "none":
        return None, "ocr_disabled"
    
    # iOS Shortcuts backend
    if backend == "shortcut":
        try:
            import shortcuts  # type: ignore
            shortcut_name = ocr_config.get("shortcut_name", "FolderExtractor OCR")
            result = shortcuts.run(shortcut_name, input=str(path))
            if isinstance(result, str) and result.strip():
                return result, "ocr_shortcut"
        except Exception:
            pass
        return None, "ocr_shortcut_error"
    
    # Tesseract OCR
    if backend == "tesseract":
        try:
            from PIL import Image
            import pytesseract
            image = Image.open(path)
            text = pytesseract.image_to_string(image)
            if text.strip():
                return text, "ocr_tesseract"
        except ImportError:
            return None, "ocr_tesseract_not_installed"
        except Exception as e:
            return None, f"ocr_error: {e}"
    
    return None, "ocr_unsupported_backend"


def extract_office(path: Path) -> Tuple[Optional[str], str]:
    """
    Extract text from Office documents (.docx, .pptx, .xlsx).
    Returns (text, method)
    """
    ext = path.suffix.lower()
    
    # Word documents
    if ext in {".docx", ".doc"}:
        try:
            import docx
            doc = docx.Document(path)
            paragraphs = [p.text for p in doc.paragraphs if p.text.strip()]
            if paragraphs:
                return "\n\n".join(paragraphs), "office_docx"
        except ImportError:
            return None, "office_docx_not_installed"
        except Exception as e:
            return None, f"office_error: {e}"
    
    # PowerPoint documents
    if ext in {".pptx", ".ppt"}:
        try:
            from pptx import Presentation
            prs = Presentation(path)
            text_parts = []
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        text_parts.append(shape.text)
            if text_parts:
                return "\n\n".join(text_parts), "office_pptx"
        except ImportError:
            return None, "office_pptx_not_installed"
        except Exception as e:
            return None, f"office_error: {e}"
    
    # Excel documents
    if ext in {".xlsx", ".xls"}:
        try:
            import openpyxl
            wb = openpyxl.load_workbook(path, data_only=True)
            text_parts = []
            for sheet in wb.worksheets:
                text_parts.append(f"Sheet: {sheet.title}")
                for row in sheet.iter_rows(values_only=True):
                    row_text = "\t".join(str(cell) if cell is not None else "" for cell in row)
                    if row_text.strip():
                        text_parts.append(row_text)
            if text_parts:
                return "\n\n".join(text_parts), "office_xlsx"
        except ImportError:
            return None, "office_xlsx_not_installed"
        except Exception as e:
            return None, f"office_error: {e}"
    
    return None, "office_unsupported"


# ========= Main Extraction Logic =========

def gather_files(root: Path) -> List[Tuple[Path, Path]]:
    """Gather all files from directory."""
    files: List[Tuple[Path, Path]] = []
    
    for dirpath, dirnames, filenames in os.walk(root):
        d = Path(dirpath)
        
        # Filter directories (keep .github as it's often important)
        dirnames[:] = [
            dn for dn in dirnames
            if dn not in IGNORE_DIR_NAMES and (not dn.startswith(".") or dn == ".github")
        ]
        
        for fn in filenames:
            p = d / fn
            if not p.is_file():
                continue
            if any(p.name.endswith(suf) for suf in IGNORE_FILE_SUFFIXES):
                continue
            
            rel = p.relative_to(root)
            files.append((p, rel))
    
    files.sort(key=lambda t: str(t[1]).lower())
    return files


def extract_file_content(
    path: Path,
    file_type: str,
    max_file_bytes: int,
    ocr_config: Dict[str, Any]
) -> Tuple[Optional[str], str]:
    """
    Extract content from file based on type.
    Returns (content, method)
    """
    if file_type == "text":
        return extract_text(path, max_file_bytes)
    elif file_type == "pdf":
        return extract_pdf(path)
    elif file_type == "image":
        return extract_image_ocr(path, ocr_config)
    elif file_type == "office":
        return extract_office(path)
    elif file_type == "binary":
        return None, "binary_skipped"
    else:
        return None, "unknown_type"


def write_file_section(
    out: io.TextIOBase,
    rel_path: Path,
    file_type: str,
    size: int,
    md5: str,
    content: Optional[str],
    method: str
) -> None:
    """Write a file section to output."""
    out.write(f"## File: {rel_path}\n\n")
    out.write(f"- **Type:** {file_type}\n")
    out.write(f"- **Size:** {human_size(size)}\n")
    out.write(f"- **Method:** {method}\n")
    if md5:
        out.write(f"- **MD5:** `{md5}`\n")
    out.write("\n")
    
    if content:
        lang = language_for(rel_path) if file_type == "text" else "text"
        out.write(f"```{lang}\n")
        out.write(content)
        if not content.endswith("\n"):
            out.write("\n")
        out.write("```\n\n")
    else:
        if method.startswith("error"):
            out.write(f"> ❌ {method}\n\n")
        elif "not_installed" in method:
            out.write(f"> ⚠️  Library not installed: {method}\n\n")
        elif method in {"binary_skipped", "ocr_disabled"}:
            out.write(f"> ℹ️  Not processed: {method}\n\n")
        else:
            out.write(f"> ⚠️  No content extracted: {method}\n\n")


def generate_extraction(
    root: Path,
    output_path: Path,
    max_file_bytes: int,
    max_output_bytes: int
) -> None:
    """Generate folder extraction."""
    cfg = load_config()
    ocr_config = cfg.get("ocr", {})
    
    # Gather files
    files = gather_files(root)
    
    if not files:
        print(f"⚠️  No files found in {root}")
        return
    
    # Statistics
    stats = {
        "total_files": len(files),
        "total_bytes": 0,
        "types": {},
        "methods": {},
    }
    
    # Generate content
    from io import StringIO
    buffer = StringIO()
    
    # Header
    buffer.write(f"# Folder Extractor Report: {root.name}\n\n")
    buffer.write(f"- **Source:** `{root}`\n")
    buffer.write(f"- **Generated:** {datetime.now().strftime('%Y-%m-%d %H:%M')}\n")
    buffer.write(f"- **Files:** {len(files)}\n\n")
    
    # Process each file
    for abs_path, rel_path in files:
        try:
            size = abs_path.stat().st_size
            stats["total_bytes"] += size
        except Exception:
            size = 0
        
        md5 = file_md5(abs_path)
        file_type = detect_file_type(abs_path)
        
        stats["types"][file_type] = stats["types"].get(file_type, 0) + 1
        
        # Extract content
        content, method = extract_file_content(
            abs_path, file_type, max_file_bytes, ocr_config
        )
        
        stats["methods"][method] = stats["methods"].get(method, 0) + 1
        
        # Write section
        write_file_section(buffer, rel_path, file_type, size, md5, content, method)
    
    # Summary
    buffer.write("---\n\n")
    buffer.write("## Summary\n\n")
    buffer.write(f"- **Total Files:** {stats['total_files']}\n")
    buffer.write(f"- **Total Size:** {human_size(stats['total_bytes'])}\n\n")
    buffer.write("**File Types:**\n")
    for ftype, count in sorted(stats["types"].items()):
        buffer.write(f"- {ftype}: {count}\n")
    buffer.write("\n**Extraction Methods:**\n")
    for method, count in sorted(stats["methods"].items()):
        buffer.write(f"- {method}: {count}\n")
    buffer.write("\n")
    
    content = buffer.getvalue()
    
    # Split if needed
    if max_output_bytes > 0 and len(content.encode(ENCODING)) > max_output_bytes:
        # Split into parts
        lines = content.split("\n")
        parts = []
        current_part = []
        current_size = 0
        
        for line in lines:
            line_bytes = len(line.encode(ENCODING)) + 1
            if current_size + line_bytes > max_output_bytes and current_part:
                parts.append("\n".join(current_part))
                current_part = [line]
                current_size = line_bytes
            else:
                current_part.append(line)
                current_size += line_bytes
        
        if current_part:
            parts.append("\n".join(current_part))
        
        # Write parts with proper headers
        stem = output_path.stem
        parent = output_path.parent
        suffix = output_path.suffix
        
        for i, part in enumerate(parts, 1):
            part_path = parent / f"{stem}_part{i}{suffix}"
            
            # Add part header
            part_content = f"# Folder Extractor Report: {root.name} (Part {i}/{len(parts)})\n\n"
            part_content += part
            
            part_path.write_text(part_content, encoding=ENCODING)
            print(f"✅ Part {i}/{len(parts)} created: {part_path}")
    else:
        # Single file
        output_path.write_text(content, encoding=ENCODING)
        print(f"✅ Extraction complete: {output_path}")


# ========= CLI =========

def parse_args(argv: List[str]) -> Dict[str, Any]:
    """Parse command line arguments."""
    args = {
        "root": None,
        "out": None,
        "max_file_bytes": DEFAULT_MAX_FILE_BYTES,
        "max_output_bytes": DEFAULT_MAX_OUTPUT_BYTES,
    }
    
    it = iter(argv)
    for token in it:
        if token in ("--root", "-r"):
            args["root"] = next(it, None)
        elif token in ("--out", "-o"):
            args["out"] = next(it, None)
        elif token == "--max-file-bytes":
            try:
                args["max_file_bytes"] = int(next(it, str(DEFAULT_MAX_FILE_BYTES)))
            except ValueError:
                pass
        elif token == "--max-output-bytes":
            try:
                args["max_output_bytes"] = int(next(it, str(DEFAULT_MAX_OUTPUT_BYTES)))
            except ValueError:
                pass
        elif token in ("--help", "-h"):
            print_help()
            sys.exit(0)
        elif not token.startswith("-") and args["root"] is None:
            args["root"] = token
    
    return args


def print_help() -> None:
    """Print help message."""
    print("""
folder-extractor – Universal Folder to Text Converter

Usage:
  folder-extractor [options] [root]
  
Options:
  --root, -r PATH          Source folder (default: current directory)
  --out, -o PATH           Output file path (default: auto-generated)
  --max-file-bytes BYTES   Max bytes to read per file (default: 10MB)
  --max-output-bytes BYTES Max bytes per output file, enables splitting (default: 5MB)
  --help, -h               Show this help message

Supported File Types:
  - Text files: direct read
  - PDFs: text extraction (requires PyPDF2 or pdfplumber)
  - Images: OCR (requires pytesseract or iOS Shortcuts)
  - Office: .docx, .pptx, .xlsx (requires python-docx, python-pptx, openpyxl)

Configuration:
  Create ~/.config/folder-extractor/config.toml for OCR settings:
  
  # For Tesseract OCR:
  [ocr]
  backend = "tesseract"
  
  # For iOS Shortcuts OCR:
  [ocr]
  backend = "shortcut"
  shortcut_name = "FolderExtractor OCR"

Examples:
  folder-extractor --root /path/to/folder --out folder_dump.md
  folder-extractor . --max-output-bytes 10000000
""")


def main() -> None:
    """Main entry point."""
    args = parse_args(sys.argv[1:])
    
    # Determine root
    root_str = args["root"] or "."
    root = Path(root_str).expanduser().resolve()
    
    if not root.is_dir():
        print(f"❌ Error: Not a directory: {root}", file=sys.stderr)
        sys.exit(1)
    
    # Determine output path
    if args["out"]:
        output_path = Path(args["out"]).expanduser().resolve()
    else:
        timestamp = datetime.now().strftime("%Y%m%d-%H%M")
        output_path = root.parent / f"{root.name}_extraction_{timestamp}.md"
    
    # Generate extraction
    try:
        generate_extraction(
            root,
            output_path,
            args["max_file_bytes"],
            args["max_output_bytes"]
        )
    except Exception as e:
        print(f"❌ Error during extraction: {e}", file=sys.stderr)
        import traceback
        traceback.print_exc()
        sys.exit(1)


if __name__ == "__main__":
    main()
